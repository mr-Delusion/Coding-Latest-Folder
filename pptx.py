# Redefine the full list of slides to avoid the NameError

from pptx import Presentation

prs = Presentation()

slides_content = [
    {
        "title": "Cleaning the Environment",
        "content": "A Step Toward a Healthier Future\nName:\nGrade:\nSubject: CSL Project\n[Space for a relevant image of clean environment or community cleanup]"
    },
    {
        "title": "Introduction",
        "content": "• The environment is our shared home.\n• A clean environment is essential for healthy living.\n• Everyone has a role to play in keeping it clean.\n[Space for picture: Children cleaning a playground or adults sweeping streets]"
    },
    {
        "title": "Why Cleaning the Environment Matters",
        "content": "• Reduces spread of diseases.\n• Improves quality of life.\n• Creates a beautiful and safe community.\n• Helps protect animals and plants.\n[Space for picture: Before and after cleanup area]"
    },
    {
        "title": "Common Environmental Problems in African Neighbourhoods",
        "content": "• Unsorted waste and garbage piles.\n• Pollution in rivers and streams.\n• Burning plastic and open dumping.\n• Blocked drainage leading to flooding.\n[Space for picture: Unsorted waste or open dumping site]"
    },
    {
        "title": "What is Waste Sorting?",
        "content": "• Separating waste into different categories:\n  - Organic (food, garden waste)\n  - Recyclable (plastic, paper, metal)\n  - Non-recyclable (used diapers, dirty packaging)\n• Helps in better waste disposal and recycling.\n[Space for picture: Colored waste bins or labeled sorting bins]"
    },
    {
        "title": "Benefits of Waste Sorting in African Neighbourhoods",
        "content": "• Makes recycling possible and easier.\n• Reduces amount of waste in landfills.\n• Creates job opportunities in recycling and waste management.\n• Keeps streets clean and prevents diseases.\n• Builds awareness and responsibility in the community.\n[Space for picture: Community waste sorting activity or recycling center]"
    },
    {
        "title": "How to Start Waste Sorting at Home",
        "content": "• Use separate bins for different types of waste.\n• Educate family members about what goes where.\n• Collaborate with neighbors and local leaders.\n• Encourage local councils to support recycling efforts.\n[Space for picture: Family sorting waste or children using labeled bins]"
    },
    {
        "title": "Community Involvement",
        "content": "• Organize cleanup events.\n• Set up waste collection points.\n• Share knowledge through schools and churches.\n• Support local waste recyclers and businesses.\n[Space for picture: Community cleanup day or group volunteers]"
    },
    {
        "title": "Conclusion",
        "content": "• Cleaning the environment starts with you.\n• Sorting waste is a simple but powerful action.\n• Together, we can make our neighbourhoods healthier and cleaner.\n• Let’s be proud of a clean Africa!\n[Space for picture: Smiling children or clean community scene]"
    },
    {
        "title": "References (Optional)",
        "content": "• Add any books, websites, or people consulted.\n[Leave this blank if not required]"
    },
    {
        "title": "Thank You!",
        "content": "[Space for image: Smiling group or thank-you banner]\nEncourage questions or discussion if presenting live."
    }
]

# Create the slides
for slide_info in slides_content:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = slide_info["title"]
    content.text = slide_info["content"]

# Save the presentation
pptx_path = "/mnt/data/Cleaning_Environment_CSL_Project.pptx"
prs.save(pptx_path)

pptx_path
